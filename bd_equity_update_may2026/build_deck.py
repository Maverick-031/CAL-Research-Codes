"""
CAL Bangladesh Research — Bangladesh Equity Market Update
Deliverable for Commercial Bank of Ceylon | May 2026

SUPERSEDED — this is the original 8-slide generator (trailing P/E 10.1x working set).
The shipped deliverable was later hand-edited (P/E refreshed to 8.5x, disclaimer added)
and is post-processed by add_narrative_slide.py. This script is kept for reference only
and now writes to a clearly-marked legacy filename so it cannot overwrite the deliverable.

Brand: C6198D (magenta), 2F3691 (deep blue), 01A7E1 (cyan). Font: DM Sans.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------- BRAND ----------
DEEP_BLUE = RGBColor(0x2F, 0x36, 0x91)
MAGENTA   = RGBColor(0xC6, 0x19, 0x8D)
CYAN      = RGBColor(0x01, 0xA7, 0xE1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1A, 0x1A, 0x1A)
INK_MID   = RGBColor(0x55, 0x55, 0x55)
INK_LIGHT = RGBColor(0x8A, 0x8A, 0x8A)
LINE      = RGBColor(0xE6, 0xE6, 0xE6)
SOFT_BG   = RGBColor(0xF7, 0xF8, 0xFB)

FONT = "DM Sans"

# 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]

# ----------------------- helpers ------------------------------

def add_slide():
    s = prs.slides.add_slide(BLANK)
    # white background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, font=FONT, size=11, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of paragraphs, each a list of (text, dict) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, opts in para:
            r = p.add_run()
            r.text = txt
            r.font.name = opts.get("font", FONT)
            r.font.size = Pt(opts.get("size", 11))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", INK)
    return tb

def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_round_rect(slide, x, y, w, h, fill, line=None, radius=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    # adjust radius
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_dot(slide, cx, cy, d, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def header_bar(slide):
    """Subtle top-of-slide brand row — text only (no decorative bars)."""
    add_text(slide, 0.55, 0.32, 8.0, 0.3,
             "CAL BANGLADESH RESEARCH",
             size=9, bold=True, color=DEEP_BLUE)
    add_text(slide, 5.0, 0.32, 8.0, 0.3,
             "EQUITY MARKET UPDATE  ·  MAY 2026",
             size=9, color=INK_LIGHT)
    # three-color motif dots
    add_dot(slide, 12.40, 0.43, 0.12, MAGENTA)
    add_dot(slide, 12.60, 0.43, 0.12, DEEP_BLUE)
    add_dot(slide, 12.80, 0.43, 0.12, CYAN)

def footer(slide, page_num, page_total, sources):
    add_text(slide, 0.55, 7.15, 10.0, 0.25,
             f"Sources: {sources}",
             size=7.5, italic=True, color=INK_LIGHT)
    add_text(slide, 11.5, 7.15, 1.3, 0.25,
             f"{page_num} / {page_total}",
             size=8, color=INK_LIGHT, align=PP_ALIGN.RIGHT)

def title_block(slide, title, subtitle):
    """Action title + descriptive subtitle (Minto)."""
    add_text(slide, 0.55, 0.80, 12.25, 1.0, title,
             size=22, bold=True, color=DEEP_BLUE, line_spacing=1.15)
    add_text(slide, 0.55, 1.78, 12.25, 0.45, subtitle,
             size=12.5, italic=True, color=INK_MID, line_spacing=1.2)

PAGE_TOTAL = 8

# ============================================================
# SLIDE 1 — COVER
# ============================================================
def slide_cover():
    s = add_slide()

    # left vertical accent — a thick magenta bar is one strong motif element on cover
    add_rect(s, 0, 0, 0.30, 7.5, MAGENTA)

    # CAL identity
    add_text(s, 0.85, 0.65, 8.0, 0.4, "CAL BANGLADESH RESEARCH",
             size=11, bold=True, color=DEEP_BLUE)
    add_text(s, 0.85, 1.05, 8.0, 0.3, "Capital Alliance Securities (Bangladesh)",
             size=9.5, italic=True, color=INK_LIGHT)

    # client tag
    add_runs(s, 0.85, 1.55, 11.0, 0.45,
             [[("PREPARED FOR  ", {"size": 9.5, "bold": True, "color": INK_LIGHT}),
               ("Commercial Bank of Ceylon PLC", {"size": 11, "bold": True, "color": MAGENTA})]])

    # main title — large
    add_text(s, 0.85, 2.45, 12.0, 1.4,
             "Bangladesh Equities:",
             size=46, bold=True, color=DEEP_BLUE, line_spacing=1.0)
    add_text(s, 0.85, 3.30, 12.0, 1.4,
             "The Re-Rating Setup",
             size=46, bold=True, color=MAGENTA, line_spacing=1.0)

    # subtitle / standfirst
    add_text(s, 0.85, 4.60, 11.5, 1.5,
             "Why a ~10x trailing P/E, a rebuilt regulatory rulebook and an incoming "
             "leadership reset at BSEC create an asymmetric entry into one of Asia's "
             "most under-owned markets — into the BNP government's first budget.",
             size=15, italic=True, color=INK_MID, line_spacing=1.35)

    # three-color motif dots row
    add_dot(s, 0.95, 6.05, 0.16, MAGENTA)
    add_dot(s, 1.25, 6.05, 0.16, DEEP_BLUE)
    add_dot(s, 1.55, 6.05, 0.16, CYAN)

    # date / authoring line
    add_text(s, 0.85, 6.40, 11.5, 0.3,
             "EQUITY MARKET UPDATE   ·   19 MAY 2026",
             size=10, bold=True, color=DEEP_BLUE)
    add_text(s, 0.85, 6.75, 11.5, 0.3,
             "Authored by the CAL Bangladesh Research team. For institutional client use only.",
             size=9, color=INK_LIGHT)

slide_cover()

# ============================================================
# SLIDE 2 — EXECUTIVE SUMMARY (Minto governing thought)
# ============================================================
def slide_summary():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "Bangladesh equities offer asymmetric upside at ~10x P/E as reforms compound into the BNP government's first budget",
        "Our governing thesis ahead of FY27 — three pillars support the call"
    )

    # governing-thought callout band
    add_round_rect(s, 0.55, 2.45, 12.25, 1.05, SOFT_BG, radius=0.10)
    add_rect(s, 0.55, 2.45, 0.10, 1.05, MAGENTA)
    add_runs(s, 0.85, 2.58, 11.85, 0.85,
        [[("Governing thought   ", {"size": 9.5, "bold": True, "color": MAGENTA})],
         [("Bangladesh equities trade at a ", {"size": 13.5, "color": INK}),
          ("~17% discount to their 10-year average ", {"size": 13.5, "bold": True, "color": DEEP_BLUE}),
          ("and a ", {"size": 13.5, "color": INK}),
          ("~44% discount to Vietnam, ", {"size": 13.5, "bold": True, "color": DEEP_BLUE}),
          ("even as a multi-year regulatory reset and a stabilising macro create the conditions for a re-rating.",
           {"size": 13.5, "color": INK})]],
        line_spacing=1.25
    )

    # three supporting pillars (cards)
    top_y = 3.85
    card_h = 2.85
    card_w = 3.92
    gap = 0.18
    starts = [0.55, 0.55 + card_w + gap, 0.55 + 2*(card_w + gap)]

    pillars = [
        {
            "tag": "01  ·  VALUATION",
            "color": MAGENTA,
            "headline": "Cheap on every metric",
            "body": [
                ("~10x", "Trailing P/E vs 12.1x 10-yr avg & 18.0x Vietnam"),
                ("-29%", "DSEX still below its Sep-2021 all-time high"),
                ("~19.5%", "Market cap / GDP — frontier-leading penetration room"),
            ]
        },
        {
            "tag": "02  ·  HEADWINDS",
            "color": DEEP_BLUE,
            "headline": "Risks are dated and priced",
            "body": [
                ("10.00%", "Repo rate held — peak monetary tightness in view"),
                ("3.9%", "WB FY26 GDP cut (Apr-26) — bad news in the price"),
                ("30.6%", "System NPL — recognition shock now visible"),
            ]
        },
        {
            "tag": "03  ·  TAILWINDS",
            "color": CYAN,
            "headline": "Reforms compound",
            "body": [
                ("FY27 Budget", "BNP govt's first budget on 11 Jun-26 — tax-relief asks queued"),
                ("BSEC reset", "New chairman appointment imminent — institutional, not personal"),
                ("CCBL + T+1", "MSCI accessibility plumbing in build-out"),
            ]
        },
    ]

    for x, p in zip(starts, pillars):
        add_round_rect(s, x, top_y, card_w, card_h, WHITE, line=LINE, radius=0.05)
        # top color strip
        add_rect(s, x, top_y, card_w, 0.10, p["color"])
        add_text(s, x + 0.25, top_y + 0.22, card_w - 0.5, 0.3, p["tag"],
                 size=8.5, bold=True, color=p["color"])
        add_text(s, x + 0.25, top_y + 0.55, card_w - 0.5, 0.45, p["headline"],
                 size=15.5, bold=True, color=DEEP_BLUE)
        # body stat rows
        ry = top_y + 1.10
        for stat, label in p["body"]:
            add_text(s, x + 0.25, ry, 1.55, 0.32, stat,
                     size=13, bold=True, color=p["color"])
            add_text(s, x + 1.85, ry + 0.04, card_w - 2.05, 0.55, label,
                     size=9, color=INK_MID, line_spacing=1.2)
            ry += 0.55

    footer(s, 2, PAGE_TOTAL, "CAL Bangladesh Research synthesis; CEIC; DSE; TBS; IMF (Jan-26 Article IV); World Bank (Apr-26 BDU); BSEC")

slide_summary()

# ============================================================
# SLIDE 3 — VALUATION (Undervalued)
# ============================================================
def slide_valuation():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "DSEX trades at a 17% discount to history and a 44% discount to Vietnam — every lens flashes value",
        "P/E, drawdown and market-cap-to-GDP all point to a market priced for bad news"
    )

    # ---- LEFT: P/E peer bars ----
    panel_y = 2.40
    panel_h = 4.50
    add_text(s, 0.55, panel_y, 5.6, 0.35,
             "Trailing P/E — DSEX vs frontier peers",
             size=12, bold=True, color=DEEP_BLUE)
    add_text(s, 0.55, panel_y + 0.32, 5.6, 0.28,
             "Bangladesh discounts Vietnam by ~44%",
             size=9.5, italic=True, color=INK_LIGHT)

    # bars
    bars = [
        ("Egypt (select banks)", 4.0,  CYAN),
        ("Pakistan KSE-100",     5.5,  CYAN),
        ("Bangladesh DSEX",      10.1, MAGENTA),
        ("DSEX 10-yr avg",       12.1, DEEP_BLUE),
        ("Vietnam VN-Index",     18.0, INK_LIGHT),
    ]
    bar_x0 = 2.50
    bar_max_w = 3.55
    bar_h = 0.36
    bar_gap = 0.22
    row_y = panel_y + 0.95
    max_val = max(v for _, v, _ in bars)

    for label, val, col in bars:
        add_text(s, 0.55, row_y + 0.05, 1.90, 0.28, label,
                 size=9.5, color=INK)
        w = bar_max_w * (val / max_val)
        add_rect(s, bar_x0, row_y, w, bar_h, col)
        add_text(s, bar_x0 + w + 0.08, row_y + 0.05, 0.9, 0.28,
                 f"{val:.1f}x",
                 size=10.5, bold=True, color=col)
        row_y += bar_h + bar_gap

    # divider
    add_rect(s, 6.35, 2.40, 0.012, 4.45, LINE)

    # ---- RIGHT: Three stat callouts ----
    rx = 6.65
    rw = 6.20

    # stat 1 — drawdown
    add_text(s, rx, panel_y, rw, 0.32,
             "1   DSEX drawdown from peak",
             size=10.5, bold=True, color=MAGENTA)
    add_runs(s, rx, panel_y + 0.35, rw, 0.7,
        [[("-28.6%", {"size": 38, "bold": True, "color": DEEP_BLUE})]])
    add_text(s, rx, panel_y + 1.05, rw, 0.35,
             "5,229 (12 May 2026) vs 7,329 all-time high (Sep 2021). A 41% rally off the Jan-2025 trough still leaves substantial valuation runway.",
             size=10, color=INK_MID, line_spacing=1.25)

    # stat 2 — discount to history
    add_text(s, rx, panel_y + 1.70, rw, 0.32,
             "2   Discount to 10-year average P/E",
             size=10.5, bold=True, color=CYAN)
    add_runs(s, rx, panel_y + 2.05, rw, 0.7,
        [[("-16.5%", {"size": 38, "bold": True, "color": DEEP_BLUE})]])
    add_text(s, rx, panel_y + 2.75, rw, 0.35,
             "Re-rating to the long-run mean alone implies +19.8% upside, ex-EPS growth. End-2025 low was 8.6x.",
             size=10, color=INK_MID, line_spacing=1.25)

    # stat 3 — penetration
    add_text(s, rx, panel_y + 3.40, rw, 0.32,
             "3   Market cap to GDP",
             size=10.5, bold=True, color=DEEP_BLUE)
    add_runs(s, rx, panel_y + 3.75, rw, 0.7,
        [[("~19.5%", {"size": 38, "bold": True, "color": DEEP_BLUE})]])
    add_text(s, rx, panel_y + 4.45, rw, 0.35,
             "Vietnam: 51% · EM avg: >70% · A long structural runway as capital formation deepens.",
             size=10, color=INK_MID, line_spacing=1.25)

    footer(s, 3, PAGE_TOTAL,
           "CEIC (Feb 2026); DSE; Taprobane Market Lens; CAL calculations. P/E peer values are trailing.")

slide_valuation()

# ============================================================
# SLIDE 4 — Liquidity & Flows
# ============================================================
def slide_liquidity():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "Liquidity is rebuilding off a five-year low — but foreign capital is still a net seller",
        "Turnover has doubled off the Jan-2025 trough; FPI outflows widened 7x in Jul–Oct FY26"
    )

    # ---- TOP: Turnover trend (sparkline-style bar series) ----
    top_y = 2.40
    panel_h = 2.65
    add_round_rect(s, 0.55, top_y, 7.55, panel_h, SOFT_BG, radius=0.04)
    add_text(s, 0.80, top_y + 0.20, 7.0, 0.32,
             "DSE average daily turnover, BDT crore",
             size=11.5, bold=True, color=DEEP_BLUE)
    add_text(s, 0.80, top_y + 0.50, 7.0, 0.28,
             "Off the early-2025 trough but still well below the 2021 peak",
             size=9, italic=True, color=INK_LIGHT)

    # bar series (illustrative monthly path)
    series = [
        ("Jan-25", 413),
        ("Mar-25", 510),
        ("Jun-25", 605),
        ("Sep-25", 670),
        ("Dec-25", 720),
        ("Mar-26", 820),
        ("May-26", 914),
    ]
    chart_x = 0.85
    chart_y = top_y + 1.00
    chart_w = 7.0
    chart_h = 1.40
    max_v = 1000  # ceiling
    bar_w = chart_w / (len(series) * 1.6)
    gap = bar_w * 0.6
    cur_x = chart_x
    for label, v in series:
        h = chart_h * (v / max_v)
        add_rect(s, cur_x, chart_y + chart_h - h, bar_w, h, CYAN)
        add_text(s, cur_x - 0.05, chart_y + chart_h - h - 0.30, bar_w + 0.30, 0.25,
                 f"{v}", size=8, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)
        add_text(s, cur_x - 0.05, chart_y + chart_h + 0.05, bar_w + 0.30, 0.25,
                 label, size=8, color=INK_MID, align=PP_ALIGN.CENTER)
        cur_x += bar_w + gap

    # peak reference — top-right corner of the chart panel (no overlap with x-axis labels)
    add_text(s, 5.10, top_y + 0.55, 2.85, 0.28,
             "Ref: 2021 peak >2,000 cr/day",
             size=8.5, italic=True, color=MAGENTA, align=PP_ALIGN.RIGHT)

    # ---- TOP-RIGHT: FPI flows card ----
    fx = 8.30
    fw = 4.55
    add_round_rect(s, fx, top_y, fw, panel_h, WHITE, line=LINE, radius=0.04)
    add_rect(s, fx, top_y, fw, 0.10, MAGENTA)
    add_text(s, fx + 0.25, top_y + 0.25, fw - 0.5, 0.3,
             "FOREIGN PORTFOLIO FLOWS",
             size=9, bold=True, color=MAGENTA)
    add_text(s, fx + 0.25, top_y + 0.55, fw - 0.5, 0.35,
             "Net outflows widened 7x YoY",
             size=13.5, bold=True, color=DEEP_BLUE)

    # two-stat side-by-side
    add_text(s, fx + 0.25, top_y + 1.05, fw - 0.5, 0.30,
             "Jul–Oct FY26", size=9, color=INK_LIGHT)
    add_runs(s, fx + 0.25, top_y + 1.30, fw - 0.5, 0.5,
             [[("-USD 66 mn", {"size": 22, "bold": True, "color": MAGENTA})]])

    add_text(s, fx + 0.25, top_y + 1.85, fw - 0.5, 0.30,
             "Same period FY25", size=9, color=INK_LIGHT)
    add_runs(s, fx + 0.25, top_y + 2.10, fw - 0.5, 0.5,
             [[("-USD 9 mn", {"size": 18, "bold": True, "color": INK_MID})]])

    # ---- BOTTOM: Top-5 concentration table ----
    by = 5.25
    add_text(s, 0.55, by, 7.55, 0.32,
             "Top 5 stocks ≈ 28–30% of DSE market cap — depth is concentrated",
             size=11.5, bold=True, color=DEEP_BLUE)

    rows = [
        ("Grameenphone",       "Telecom",      "42,710", "12.0%", MAGENTA),
        ("Square Pharma",      "Pharma",       "19,129", "5.4%",  CYAN),
        ("BAT Bangladesh",     "Tobacco",      "18,673", "5.2%",  CYAN),
        ("BRAC Bank",          "Banks",        ">15,000", "4.2%", CYAN),
    ]
    # header
    th_y = by + 0.45
    add_text(s, 0.55, th_y, 2.5, 0.25, "STOCK",       size=8.5, bold=True, color=INK_LIGHT)
    add_text(s, 3.10, th_y, 1.8, 0.25, "SECTOR",       size=8.5, bold=True, color=INK_LIGHT)
    add_text(s, 5.10, th_y, 1.5, 0.25, "MKT CAP (BDT cr)", size=8.5, bold=True, color=INK_LIGHT, align=PP_ALIGN.RIGHT)
    add_text(s, 6.80, th_y, 1.2, 0.25, "% OF DSE",     size=8.5, bold=True, color=INK_LIGHT, align=PP_ALIGN.RIGHT)
    # subtle divider
    add_rect(s, 0.55, th_y + 0.30, 7.55, 0.012, LINE)

    ry = th_y + 0.42
    for name, sec, mc, pct, dot in rows:
        add_dot(s, 0.62, ry + 0.13, 0.08, dot)
        add_text(s, 0.78, ry, 2.4, 0.25, name, size=10, bold=True, color=INK)
        add_text(s, 3.10, ry, 1.8, 0.25, sec, size=9.5, color=INK_MID)
        add_text(s, 5.10, ry, 1.5, 0.25, mc, size=10, color=INK, align=PP_ALIGN.RIGHT)
        add_text(s, 6.80, ry, 1.2, 0.25, pct, size=10, bold=True, color=DEEP_BLUE, align=PP_ALIGN.RIGHT)
        ry += 0.30

    # ---- BOTTOM-RIGHT: takeaway card ----
    tx = 8.30
    tw = 4.55
    th_h = 1.85
    add_round_rect(s, tx, by, tw, th_h, SOFT_BG, radius=0.04)
    add_text(s, tx + 0.25, by + 0.18, tw - 0.5, 0.30,
             "WHAT IT MEANS", size=9, bold=True, color=MAGENTA)
    add_text(s, tx + 0.25, by + 0.50, tw - 0.5, 1.40,
             "Foreign capital is still derisking around the post-election transition; "
             "a 7x worsening in outflows is the contra-indicator we want to see fade. "
             "Domestic liquidity is leading, but breadth is shallow — entry should "
             "target index heavyweights with re-rating optionality.",
             size=10, color=INK, line_spacing=1.30)

    footer(s, 4, PAGE_TOTAL,
           "DSE; CEIC; The Business Standard; The Daily Star. Turnover series illustrative based on weekly press recaps.")

slide_liquidity()

# ============================================================
# SLIDE 5 — HEADWINDS
# ============================================================
def slide_headwinds():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "Cyclical and geopolitical headwinds anchor the entry point — but most are dated and priced",
        "Four risk vectors that have largely been absorbed by valuation; we look for inflection signals into H2-26"
    )

    # 2x2 grid
    grid_y = 2.40
    grid_h = 4.50
    cw = 6.10
    ch = 2.10
    gx = 0.55
    gy = 0.20
    positions = [
        (gx, grid_y, "Monetary tightness", MAGENTA, "Peak rates choke private credit", [
            "Repo rate held at 10.00% (Jan-26 MPS) — highest in two decades",
            "Private credit growth 6.03% YoY (Feb-26) — 5-yr low vs 8.5% target",
            "Lending rates 13.5–15%; real lending rate +5.96% — punitive for SMEs",
        ]),
        (gx + cw + gy, grid_y, "Banking sector stress", DEEP_BLUE, "Recognition shock is visible — but bounded", [
            "System NPL 30.6% (Dec-25, World Bank) — likely highest globally",
            "Islami Bank: bad loans Tk 94,322 cr (+44% YoY); NPL 51%; CRAR 6.42%",
            "S Alam Group exposure drives concentration; mergers under way",
        ]),
        (gx, grid_y + ch + gy, "Geopolitics — four-front squeeze", CYAN, "Iran war + US tariff regime + India freeze", [
            "Iran-Israel war: 40 of 115 contracted 2026 LNG cargoes expected lost",
            "US 19% reciprocal tariff (Feb-26) on RMG exports (USD 8.2bn to US in 2025)",
            "India: visa freeze (Dec-25); CEPA shelved; Red Sea freight +20–50%",
        ]),
        (gx + cw + gy, grid_y + ch + gy, "Growth & IMF program", MAGENTA, "Downgrades pricing in slower convergence", [
            "FY26 GDP cut to 3.9% — World Bank BDU 8 Apr-26 (was 4.6–4.8%)",
            "IMF Article IV (26 Jan-26); program extended 6m to 27 Jan-27",
            "5th-review tranche (~USD 400m) awaits BNP govt sign-off",
        ]),
    ]

    for x, y, head, color, sub, bullets in positions:
        add_round_rect(s, x, y, cw, ch, WHITE, line=LINE, radius=0.04)
        # left color tab
        add_rect(s, x, y, 0.10, ch, color)
        add_text(s, x + 0.28, y + 0.15, cw - 0.5, 0.32, head,
                 size=13, bold=True, color=DEEP_BLUE)
        add_text(s, x + 0.28, y + 0.48, cw - 0.5, 0.28, sub,
                 size=9.5, italic=True, color=color)
        # bullets
        by_ = y + 0.85
        for b in bullets:
            add_text(s, x + 0.28, by_, 0.18, 0.25, "›",
                     size=11, bold=True, color=color)
            add_text(s, x + 0.50, by_, cw - 0.75, 0.30, b,
                     size=10, color=INK, line_spacing=1.25)
            by_ += 0.36

    footer(s, 5, PAGE_TOTAL,
           "Bangladesh Bank H1-FY26 MPS; World Bank BDU 8 Apr-26; IMF PR 26/029 & 25/369 (Art-IV / 5th-review mission); TBS; Daily Star; Al Jazeera; White House.")

slide_headwinds()

# ============================================================
# SLIDE 6 — TAILWINDS PART 1: Regulatory reset
# ============================================================
def slide_tailwinds_reg():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "BSEC's rulebook reset is institutional, not personal — and a leadership handover is now imminent",
        "~Tk 1,500 cr in fines, a rebuilt rulebook, and a market-credible successor in the frame"
    )

    # left: institutional reset + leadership-transition panel
    lx = 0.55
    ly = 2.40
    lw = 5.90
    lh = 4.50

    add_round_rect(s, lx, ly, lw, lh, SOFT_BG, radius=0.04)
    add_text(s, lx + 0.28, ly + 0.22, lw - 0.5, 0.30,
             "INSTITUTIONAL RESET",
             size=9, bold=True, color=MAGENTA)
    add_text(s, lx + 0.28, ly + 0.55, lw - 0.5, 0.45,
             "Rulebook over personalities",
             size=18, bold=True, color=DEEP_BLUE)
    add_text(s, lx + 0.28, ly + 0.98, lw - 0.5, 0.32,
             "Enforcement, margin, REIT and MF rules rebuilt 2024–25  ·  test now is continuity",
             size=10, italic=True, color=INK_MID)

    # big stat — fines (kept as the institutional anchor; not tied to any individual)
    add_text(s, lx + 0.28, ly + 1.55, lw - 0.5, 0.30,
             "Cumulative market-manipulation fines, 2024–25",
             size=10, color=INK_MID)
    add_runs(s, lx + 0.28, ly + 1.85, lw - 0.5, 0.85,
             [[("Tk ", {"size": 22, "bold": True, "color": MAGENTA}),
               ("1,500", {"size": 46, "bold": True, "color": MAGENTA}),
               (" cr", {"size": 22, "bold": True, "color": MAGENTA})]])
    add_text(s, lx + 0.28, ly + 2.78, lw - 0.5, 0.30,
             "Largest enforcement sweep in BSEC's history — recovery only ~0.35%",
             size=10, italic=True, color=INK_LIGHT)

    # leadership transition note (subtle, factual)
    add_text(s, lx + 0.28, ly + 3.20, lw - 0.5, 0.28,
             "LEADERSHIP TRANSITION",
             size=8.5, bold=True, color=DEEP_BLUE)
    notes = [
        ("Chairman transition imminent",       "Press reporting May-26: handover under way"),
        ("Age-limit bar being eased",          "Opens slate to senior market-credible names"),
        ("What we want to see",                "Non-political technocrat → FPI sentiment lift"),
    ]
    cy = ly + 3.50
    for name, desc in notes:
        add_dot(s, lx + 0.36, cy + 0.09, 0.08, MAGENTA)
        add_text(s, lx + 0.52, cy - 0.03, lw - 0.8, 0.22, name,
                 size=9.5, bold=True, color=INK)
        add_text(s, lx + 0.52, cy + 0.17, lw - 0.8, 0.22, desc,
                 size=8.5, color=INK_MID)
        cy += 0.43

    # right: reform timeline
    rx = 6.65
    rw = 6.20

    add_text(s, rx, 2.40, rw, 0.32,
             "REFORM TIMELINE",
             size=9, bold=True, color=CYAN)
    add_text(s, rx, 2.70, rw, 0.32,
             "Six dated milestones rebuilt the rulebook",
             size=14, bold=True, color=DEEP_BLUE)

    milestones = [
        ("Jul 2022",  "Floor prices imposed",       "Index-defence move trapped >60% of scrips",     INK_LIGHT),
        ("21 Jan 2024", "Floor prices lifted",      "Restored price discovery; circuit breaker reset to 10%", MAGENTA),
        ("2024",      "REIT Fund Rules approved",   "Tk 200 cr min; 90% dividend mandate",            CYAN),
        ("24 Jun 2025","MSCI flags floor progress", "BD retained Frontier; under active monitoring", DEEP_BLUE),
        ("6 Nov 2025","Margin Loan Rules 2025",     "1:1 cap; 25% single-stock; 30% sector limits",  DEEP_BLUE),
        ("May 2026",  "BSEC chairman handover",     "Transition in motion; technocrat slate widened",  MAGENTA),
    ]
    ty = 3.25
    line_x = rx + 0.95
    # vertical line down the timeline
    add_rect(s, line_x, ty + 0.10, 0.018, len(milestones)*0.58 - 0.1, LINE)
    for date, head, desc, c in milestones:
        # date
        add_text(s, rx, ty, 0.95, 0.30, date,
                 size=9, bold=True, color=INK_MID, align=PP_ALIGN.RIGHT)
        # dot
        add_dot(s, line_x + 0.009, ty + 0.18, 0.16, c)
        # text
        add_text(s, line_x + 0.30, ty - 0.02, rw - 1.25, 0.30, head,
                 size=10.5, bold=True, color=DEEP_BLUE)
        add_text(s, line_x + 0.30, ty + 0.22, rw - 1.25, 0.30, desc,
                 size=9, color=INK_MID)
        ty += 0.58

    footer(s, 6, PAGE_TOTAL,
           "BSEC; TBS; BSS; Daily Sun (May-26 chairman transition); MSCI 2025 Market Classification Review (24 Jun 2025); Daily Star.")

slide_tailwinds_reg()

# ============================================================
# SLIDE 7 — TAILWINDS PART 2: Market infrastructure
# ============================================================
def slide_tailwinds_infra():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "Plumbing upgrades — CCBL, T+1, REITs and sukuk — unlock MSCI accessibility and frontier-leading depth",
        "Six structural enablers compounding through 2026–27"
    )

    # 3x2 grid of infra cards — shrunk slightly to leave room for the bottom accessibility band
    gy = 2.40
    gx = 0.55
    cw = 4.05
    ch = 1.92
    gap = 0.13

    cards = [
        ("CCBL", "Central Counterparty Bangladesh Ltd",
         "Incorporated Jan-19; targeted go-live alongside derivatives. Slips on director-appointment delays but firmly on roadmap.",
         "Slated", MAGENTA),
        ("T+1", "Settlement cycle upgrade",
         "BSEC and DSE re-engaged BB Gov. Mansur in 2025. Standard MSCI/FTSE accessibility tick — India live since 2023.",
         "In flight", DEEP_BLUE),
        ("Sukuk", "Sovereign Islamic depth",
         "Cumulative issuance ~Tk 22,000 cr by mid-2025. CDWSP Social Impact Sukuk: 5-yr, 10.40% rental, 2.5x oversubscribed.",
         "Scaling", CYAN),
        ("REITs", "REIT Fund Rules 2024",
         "Min Tk 200 cr; sponsor ≥20%; 90% dividend mandate; 3-yr lock-in. Satellite-city scope amendment under review.",
         "Live", MAGENTA),
        ("T-Bonds on DSE", "Secondary fixed income",
         "Listed since Oct-22; 2024 DSE turnover Tk 1.31 bn (+55% YoY). Reform push to extend NSCs and corp bonds.",
         "Scaling", DEEP_BLUE),
        ("MSCI / FTSE", "Index monitoring",
         "MSCI 24-Jun-25 review: BD remains Frontier, accessibility under active monitoring; next read Nov-26.",
         "Watch", CYAN),
    ]

    for i, (tag, head, body, status, c) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = gx + col * (cw + gap)
        y = gy + row * (ch + gap)
        add_round_rect(s, x, y, cw, ch, WHITE, line=LINE, radius=0.045)
        add_rect(s, x, y, cw, 0.10, c)
        # tag chip
        add_text(s, x + 0.28, y + 0.20, cw - 0.5, 0.30, tag.upper(),
                 size=9, bold=True, color=c)
        add_text(s, x + 0.28, y + 0.50, cw - 0.5, 0.32, head,
                 size=12, bold=True, color=DEEP_BLUE)
        add_text(s, x + 0.28, y + 0.83, cw - 0.5, 0.78, body,
                 size=9, color=INK_MID, line_spacing=1.25)
        # status pill (bottom)
        pill_w = 0.85
        add_round_rect(s, x + cw - pill_w - 0.22, y + ch - 0.42, pill_w, 0.28,
                       WHITE, line=c, radius=0.40)
        add_text(s, x + cw - pill_w - 0.22, y + ch - 0.39, pill_w, 0.26, status,
                 size=8, bold=True, color=c, align=PP_ALIGN.CENTER)

    # bottom takeaway band — sits comfortably above the footer at 7.15
    band_y = gy + 2*ch + gap + 0.20  # 2.40 + 3.84 + 0.13 + 0.20 = 6.57
    add_round_rect(s, 0.55, band_y, 12.30, 0.50, SOFT_BG, radius=0.30)
    add_rect(s, 0.55, band_y, 0.10, 0.50, MAGENTA)
    add_runs(s, 0.85, band_y + 0.10, 12.0, 0.35,
        [[("ACCESSIBILITY DELTA   ", {"size": 9, "bold": True, "color": MAGENTA}),
          ("Each milestone independently nudges MSCI scoring; together they raise the case for ",
           {"size": 10.5, "color": INK}),
          ("Frontier ", {"size": 10.5, "bold": True, "color": DEEP_BLUE}),
          ("→ ", {"size": 10.5, "color": INK}),
          ("Emerging review consideration", {"size": 10.5, "bold": True, "color": DEEP_BLUE}),
          (" over 2027–28.", {"size": 10.5, "color": INK})]])

    footer(s, 7, PAGE_TOTAL,
           "BSEC; CCBL; CDBL; MSCI 2025 Review; Financial Express; New Age; Traders Union; BB FX Guidelines.")

slide_tailwinds_infra()

# ============================================================
# SLIDE 8 — INVESTMENT THESIS
# ============================================================
def slide_thesis():
    s = add_slide()
    header_bar(s)

    title_block(s,
        "Our view: Overweight Bangladesh selectively over 12–18 months — own sukuk for carry, large-caps for re-rating",
        "Five catalysts to watch, five risks to underwrite, and a clean entry valuation"
    )

    # ---- TOP: Recommendation strip ----
    rec_y = 2.40
    add_round_rect(s, 0.55, rec_y, 12.30, 1.05, SOFT_BG, radius=0.08)
    add_rect(s, 0.55, rec_y, 0.10, 1.05, MAGENTA)
    add_text(s, 0.90, rec_y + 0.15, 11.7, 0.30,
             "CAL RECOMMENDATION",
             size=9, bold=True, color=MAGENTA)
    add_runs(s, 0.90, rec_y + 0.45, 11.7, 0.55,
             [[("OVERWEIGHT  ", {"size": 22, "bold": True, "color": DEEP_BLUE}),
               ("Bangladesh equities  ·  12–18m horizon  ·  ", {"size": 16, "color": INK}),
               ("Selective; ", {"size": 16, "italic": True, "color": INK_MID}),
               ("quality large-caps + sovereign sukuk for carry", {"size": 16, "bold": True, "color": MAGENTA})]])

    # ---- MIDDLE: Three columns — Where to be / Catalysts / Risks ----
    mid_y = 3.70
    mid_h = 3.10
    col_w = 4.05
    col_gap = 0.13

    # Column 1 — where to be
    x1 = 0.55
    add_round_rect(s, x1, mid_y, col_w, mid_h, WHITE, line=LINE, radius=0.04)
    add_rect(s, x1, mid_y, col_w, 0.10, MAGENTA)
    add_text(s, x1 + 0.25, mid_y + 0.22, col_w - 0.5, 0.30,
             "WHERE TO BE",
             size=9, bold=True, color=MAGENTA)
    add_text(s, x1 + 0.25, mid_y + 0.55, col_w - 0.5, 0.32,
             "Positioning playbook",
             size=12.5, bold=True, color=DEEP_BLUE)
    plays = [
        ("Quality large-caps", "USD-linked or remittance-leveraged earnings"),
        ("Pharma exporters",   "Margin defence + EM-priced multiples"),
        ("Sovereign sukuk",    "~10.4% rental locks 5-yr carry"),
        ("Banks — selective",  "Top-3 CRAR names only; avoid Islamic stress"),
        ("Avoid",              "Telecoms, NBFI, high-NPL Islamic banks"),
    ]
    py = mid_y + 1.00
    for name, desc in plays:
        add_dot(s, x1 + 0.32, py + 0.13, 0.08, MAGENTA)
        add_text(s, x1 + 0.48, py - 0.02, col_w - 0.75, 0.25, name,
                 size=10, bold=True, color=INK)
        add_text(s, x1 + 0.48, py + 0.20, col_w - 0.75, 0.25, desc,
                 size=8.8, color=INK_MID)
        py += 0.42

    # Column 2 — catalysts
    x2 = x1 + col_w + col_gap
    add_round_rect(s, x2, mid_y, col_w, mid_h, WHITE, line=LINE, radius=0.04)
    add_rect(s, x2, mid_y, col_w, 0.10, CYAN)
    add_text(s, x2 + 0.25, mid_y + 0.22, col_w - 0.5, 0.30,
             "CATALYSTS — 6 TO 12 MONTHS",
             size=9, bold=True, color=CYAN)
    add_text(s, x2 + 0.25, mid_y + 0.55, col_w - 0.5, 0.32,
             "Triggers for the re-rating",
             size=12.5, bold=True, color=DEEP_BLUE)
    cats = [
        ("1", "FY27 Budget (11 Jun-26) — CGT & IMF signal"),
        ("2", "New BSEC chairman + first 100-day signal"),
        ("3", "CCBL go-live + T+1 rollout date"),
        ("4", "MSCI Nov-26 semi-annual review"),
        ("5", "CPI <7% → BB rate-cut window"),
    ]
    py = mid_y + 1.00
    for num, desc in cats:
        # numbered circle
        add_dot(s, x2 + 0.35, py + 0.13, 0.26, CYAN)
        add_text(s, x2 + 0.23, py + 0.00, 0.30, 0.30, num,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x2 + 0.65, py + 0.06, col_w - 0.85, 0.30, desc,
                 size=10, color=INK, line_spacing=1.20)
        py += 0.42

    # Column 3 — risks
    x3 = x2 + col_w + col_gap
    add_round_rect(s, x3, mid_y, col_w, mid_h, WHITE, line=LINE, radius=0.04)
    add_rect(s, x3, mid_y, col_w, 0.10, DEEP_BLUE)
    add_text(s, x3 + 0.25, mid_y + 0.22, col_w - 0.5, 0.30,
             "RISKS TO UNDERWRITE",
             size=9, bold=True, color=DEEP_BLUE)
    add_text(s, x3 + 0.25, mid_y + 0.55, col_w - 0.5, 0.32,
             "What can break the thesis",
             size=12.5, bold=True, color=DEEP_BLUE)
    risks = [
        ("!", "IMF 5th-tranche delay / successor lapse"),
        ("!", "Iran war: LNG / power outages persist"),
        ("!", "NPL recognition wave hits bank EPS"),
        ("!", "Disorderly BSEC chairman handover"),
        ("!", "US tariff renegotiation in 2026"),
    ]
    py = mid_y + 1.00
    for mark, desc in risks:
        add_dot(s, x3 + 0.35, py + 0.13, 0.26, DEEP_BLUE)
        add_text(s, x3 + 0.23, py + 0.00, 0.30, 0.30, mark,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x3 + 0.65, py + 0.06, col_w - 0.85, 0.30, desc,
                 size=10, color=INK, line_spacing=1.20)
        py += 0.42

    footer(s, 8, PAGE_TOTAL,
           "CAL Bangladesh Research view. Sources: CEIC; DSE; BSEC; IMF (Jan-26 Art-IV); World Bank (Apr-26 BDU); TBS; Daily Star; Daily Sun; MSCI. See data_and_sources.md for full source set.")

slide_thesis()

# ============================================================
# SAVE
# ============================================================
import os
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(SCRIPT_DIR, "output")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "legacy_CAL_BD_Equity_Update_v1_8slide.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
