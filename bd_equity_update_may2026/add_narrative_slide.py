"""
Post-process the latest CBC deck (no original build script available):
  1. Insert a dedicated narrative slide as new Slide 2 (Ahmed's storyline arc).
  2. Thread the arc into the cover standfirst.
  3. Echo the '6-month easing-cycle delay' point on the headwinds slide.
  4. Renumber footer page numbers after the insert.

Edits the .pptx directly via python-pptx. Reads from a pristine copy so re-runs are clean.
"""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(SCRIPT_DIR, "output")
PRISTINE = os.path.join(SCRIPT_DIR, "CAL_BD_Equity_Update_CBC_May2026_source.pptx")
OUT = os.path.join(OUT_DIR, "CAL_BD_Equity_Update_CBC_May2026.pptx")
LOGO = os.path.join(SCRIPT_DIR, "assets", "cal_logo.png")

# ---------- BRAND ----------
DEEP_BLUE = RGBColor(0x2F, 0x36, 0x91)
MAGENTA   = RGBColor(0xC6, 0x19, 0x8D)
CYAN      = RGBColor(0x01, 0xA7, 0xE1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1A, 0x1A, 0x1A)
INK_MID   = RGBColor(0x55, 0x55, 0x55)
INK_LIGHT = RGBColor(0x8A, 0x8A, 0x8A)
LINE      = RGBColor(0xE6, 0xE6, 0xE6)
SOFT_BG   = RGBColor(0xF7, 0xF8, 0xFB)
FONT = "DM Sans"

prs = Presentation(PRISTINE)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = None
for lay in prs.slide_layouts:
    if lay.name == "Blank":
        BLANK = lay
        break
if BLANK is None:
    BLANK = prs.slide_layouts[6]

# ----------------------- helpers ------------------------------
def add_text(slide, x, y, w, h, text, *, font=FONT, size=11, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, opts in para:
            r = p.add_run()
            r.text = txt
            r.font.name = opts.get("font", FONT)
            r.font.size = Pt(opts.get("size", 11))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", INK)
    return tb

def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_round_rect(slide, x, y, w, h, fill, line=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_dot(slide, cx, cy, d, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

def replace_textbox_text(shape, new_text, *, size=None, italic=None, bold=None, color=None):
    """Rewrite a textbox to a single run, inheriting formatting from its first run."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    src = p0.runs[0] if p0.runs else None
    f_name = src.font.name if src and src.font.name else FONT
    f_size = src.font.size if src and src.font.size else None
    f_italic = src.font.italic if src else None
    f_bold = src.font.bold if src else None
    try:
        f_color = src.font.color.rgb if src and src.font.color and src.font.color.type is not None else None
    except Exception:
        f_color = None
    # clear all paragraphs except first, and clear first's runs
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    r = p0.add_run()
    r.text = new_text
    r.font.name = f_name
    r.font.size = Pt(size) if size is not None else f_size
    r.font.italic = italic if italic is not None else f_italic
    r.font.bold = bold if bold is not None else f_bold
    rgb = color if color is not None else f_color
    if rgb is not None:
        r.font.color.rgb = rgb

# ============================================================
# 1. EDIT COVER STANDFIRST (slide index 0)
# ============================================================
cover = prs.slides[0]
for sh in cover.shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("Why "):
        replace_textbox_text(
            sh,
            "Why the post-election re-rating is delayed, not derailed: a Middle East war pushed "
            "the easing cycle back ~6 months, leaving Bangladesh equities at 8.5x P/E — an "
            "asymmetric entry into one of Asia's most under-owned markets, into the BNP "
            "government's first budget.",
            size=15, italic=True, color=INK_MID)
        break

# ============================================================
# 2. EDIT HEADWINDS SLIDE — echo the 6-month delay (slide index 4)
# ============================================================
hw = prs.slides[4]
for sh in hw.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t == "Peak rates choke private credit":
        replace_textbox_text(sh, "Energy shock has reset the easing cycle")
    elif t.startswith("Lending rates 13.5"):
        replace_textbox_text(sh, "ME-war inflation spike pushed the first rate cut back ≥6 months")

# ============================================================
# 3. BUILD THE NARRATIVE SLIDE (appended, then moved to index 1)
# ============================================================
s = prs.slides.add_slide(BLANK)
# white background
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.line.fill.background(); bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
bg.shadow.inherit = False

# header: logo + header text + dots (match content slides)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.47), Inches(0.14), Inches(1.17), Inches(0.59))
add_text(s, 5.00, 0.32, 2.78, 0.17, "EQUITY MARKET UPDATE  ·  MAY 2026",
         size=9, color=INK_LIGHT)
add_dot(s, 12.40, 0.43, 0.12, MAGENTA)
add_dot(s, 12.60, 0.43, 0.12, DEEP_BLUE)
add_dot(s, 12.80, 0.43, 0.12, CYAN)

# title + subtitle
add_text(s, 0.55, 0.80, 12.25, 0.83,
         "The re-rating was set to fire in 2026 — the Middle East war reset the clock, not the case",
         size=22, bold=True, color=DEEP_BLUE, line_spacing=1.12)
add_text(s, 0.55, 1.78, 12.25, 0.45,
         "The thesis in four moves: a primed catalyst, an exogenous shock, a cheaper entry, "
         "and reforms that compound while we wait",
         size=12.5, italic=True, color=INK_MID)

# four cards
top_y = 2.45
card_h = 3.50
card_w = 2.88
gap = 0.147
xs = [0.55 + i * (card_w + gap) for i in range(4)]

cards = [
    {"tag": "01  ·  THE SETUP", "color": CYAN, "headline": "A rally was primed",
     "bul": ["Post-election stability since Feb-26",
             "Inflation rolling over from its 2024–25 peak",
             "Market positioned for the first rate cut"],
     "stat": "Feb-26", "statlbl": "BNP government takes office"},
    {"tag": "02  ·  THE SHOCK", "color": MAGENTA, "headline": "The war reset the clock",
     "bul": ["Iran–Israel war spiked energy & LNG costs",
             "Inflation & rate path knocked off track",
             "Delayed, not derailed"],
     "stat": "≥6 mo", "statlbl": "easing cycle pushed back"},
    {"tag": "03  ·  THE VALUATION", "color": DEEP_BLUE, "headline": "Cheaper than at the setup",
     "bul": ["8.5x trailing P/E — priced for bad news",
             "~30% below 10-yr avg; ~48% below Vietnam",
             "End-2025 low of 8.6x already revisited"],
     "stat": "+42%", "statlbl": "re-rating upside, ex-EPS growth"},
    {"tag": "04  ·  THE CHANGES", "color": CYAN, "headline": "Reforms keep compounding",
     "bul": ["BSEC rulebook reset + new chairman",
             "CCBL / T+1 → MSCI accessibility track",
             "Sukuk carry locks income while you wait"],
     "stat": "10.4%", "statlbl": "sovereign sukuk rental"},
]

for x, c in zip(xs, cards):
    add_round_rect(s, x, top_y, card_w, card_h, WHITE, line=LINE, radius=0.05)
    add_rect(s, x, top_y, card_w, 0.10, c["color"])
    add_text(s, x + 0.25, top_y + 0.22, card_w - 0.45, 0.28, c["tag"],
             size=8.5, bold=True, color=c["color"])
    add_text(s, x + 0.25, top_y + 0.52, card_w - 0.45, 0.55, c["headline"],
             size=14, bold=True, color=DEEP_BLUE, line_spacing=1.05)
    by = top_y + 1.18
    for b in c["bul"]:
        add_text(s, x + 0.25, by, 0.16, 0.24, "›", size=10.5, bold=True, color=c["color"])
        add_text(s, x + 0.45, by, card_w - 0.70, 0.55, b, size=9, color=INK_MID, line_spacing=1.12)
        by += 0.42
    # key stat anchor
    add_rect(s, x + 0.25, top_y + 2.66, card_w - 0.50, 0.012, LINE)
    add_text(s, x + 0.25, top_y + 2.78, card_w - 0.45, 0.42, c["stat"],
             size=20, bold=True, color=c["color"])
    add_text(s, x + 0.25, top_y + 3.20, card_w - 0.45, 0.28, c["statlbl"],
             size=8.5, color=INK_LIGHT, line_spacing=1.05)

# flow chevrons in the gaps
for i in range(3):
    cx = xs[i] + card_w + gap / 2
    add_text(s, cx - 0.15, top_y + card_h / 2 - 0.18, 0.30, 0.30, "›",
             size=18, bold=True, color=INK_LIGHT, align=PP_ALIGN.CENTER)

# bottom takeaway band
band_y = 6.15
add_round_rect(s, 0.55, band_y, 12.30, 0.55, SOFT_BG, radius=0.28)
add_rect(s, 0.55, band_y, 0.10, 0.55, MAGENTA)
add_runs(s, 0.85, band_y + 0.13, 12.0, 0.35,
    [[("NET   ", {"size": 9, "bold": True, "color": MAGENTA}),
      ("The catalyst is delayed, not cancelled — and the entry is ", {"size": 10.5, "color": INK}),
      ("cheaper than when the thesis began.", {"size": 10.5, "bold": True, "color": DEEP_BLUE}),
      (" We accumulate selectively through the delay and get paid ~10.4% to wait.",
       {"size": 10.5, "color": INK})]])

# footer
add_text(s, 0.55, 7.15, 10.0, 0.25,
         "Sources: CAL Bangladesh Research synthesis; Bangladesh Bank MPS; World Bank (Apr-26 BDU); "
         "IMF (Jan-26 Art-IV); CEIC; DSE.",
         size=7.5, italic=True, color=INK_LIGHT)
add_text(s, 11.50, 7.15, 1.30, 0.15, "2", size=8, color=INK_LIGHT, align=PP_ALIGN.RIGHT)

# ---- move the new slide (currently last) to index 1 ----
sldIdLst = prs.slides._sldIdLst
sld_ids = list(sldIdLst)
new_id = sld_ids[-1]
sldIdLst.remove(new_id)
sldIdLst.insert(1, new_id)

# ============================================================
# 4. RENUMBER FOOTER PAGE NUMBERS (bottom-right numeric boxes)
# ============================================================
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.left is None or sh.top is None:
            continue
        if Emu(sh.left).inches < 11.0 or Emu(sh.top).inches < 6.9:
            continue
        txt = sh.text_frame.text.strip()
        if txt.isdigit():
            p0 = sh.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = str(i + 1)
                for extra in p0.runs[1:]:
                    extra._r.getparent().remove(extra._r)

prs.save(OUT)
print("Saved:", OUT)
print("Total slides:", len(prs.slides))
